import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import markdown
from bs4 import BeautifulSoup

def markdown_to_docx(md_text, output_path):
    doc = Document()
    html = markdown.markdown(md_text)
    soup = BeautifulSoup(html, 'html.parser')
    
    for tag in soup.find_all(['h1', 'h2', 'h3', 'p', 'li', 'table']):
        if tag.name == 'h1':
            doc.add_heading(tag.get_text(), level=0)
        elif tag.name == 'h2':
            doc.add_heading(tag.get_text(), level=1)
        elif tag.name == 'h3':
            doc.add_heading(tag.get_text(), level=2)
        elif tag.name == 'p':
            doc.add_paragraph(tag.get_text())
        elif tag.name == 'li':
            doc.add_paragraph(tag.get_text(), style='List Bullet')
        elif tag.name == 'table':
            rows = tag.find_all('tr')
            if not rows: continue
            cols = len(rows[0].find_all(['td', 'th']))
            table = doc.add_table(rows=len(rows), cols=cols)
            table.style = 'Table Grid'
            for i, row in enumerate(rows):
                for j, cell in enumerate(row.find_all(['td', 'th'])):
                    table.cell(i, j).text = cell.get_text()
    
    doc.save(output_path)

def markdown_to_pptx(md_text, output_path):
    prs = Presentation()
    html = markdown.markdown(md_text)
    soup = BeautifulSoup(html, 'html.parser')
    
    # Simple strategy: Each H1/H2 starts a new slide
    current_slide = None
    
    for tag in soup.find_all(['h1', 'h2', 'h3', 'p', 'li']):
        if tag.name in ['h1', 'h2']:
            slide_layout = prs.slide_layouts[1] # Title and Content
            current_slide = prs.slides.add_slide(slide_layout)
            current_slide.shapes.title.text = tag.get_text()
            # Placeholder for bullets
            body_shape = current_slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""
        elif current_slide:
            body_shape = current_slide.placeholders[1]
            tf = body_shape.text_frame
            p = tf.add_paragraph()
            p.text = tag.get_text()
            if tag.name == 'li':
                p.level = 1
            elif tag.name == 'h3':
                p.font.bold = True
                p.font.size = Pt(18)
    
    prs.save(output_path)
